from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from typing import List, Optional
from pydantic import BaseModel
import shutil, os, json, re, time
import fitz
import docx
from pptx import Presentation
import numpy as np
import faiss
from sentence_transformers import SentenceTransformer
import requests
from dotenv import load_dotenv

# ── Env ──────────────────────────────────────────────
load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
if not GROQ_API_KEY:
    raise ValueError("GROQ_API_KEY not found in .env file")

GROQ_URL   = "https://api.groq.com/openai/v1/chat/completions"
GROQ_MODEL = "llama-3.3-70b-versatile"

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}

# ── Offline mode — skip HuggingFace network calls ────
os.environ["TRANSFORMERS_OFFLINE"] = "1"
os.environ["HF_DATASETS_OFFLINE"]  = "1"
os.environ["HF_HUB_OFFLINE"]       = "1"

# ── App ───────────────────────────────────────────────
app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], allow_credentials=True,
    allow_methods=["*"], allow_headers=["*"],
)

# ── Storage ───────────────────────────────────────────
UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

embed_model      = SentenceTransformer("all-MiniLM-L6-v2")
faiss_index      = faiss.IndexFlatL2(384)
documents: list  = []
file_names: list = []
doc_meta: list   = []
quiz_sessions    = {}
quiz_counter     = 0

# Session TTL (seconds)
QUIZ_SESSION_TTL = 3600  # 1 hour


# ── Startup: re-index files already on disk ───────────
# Prevents the backend losing all file data on restart.
# Any file sitting in the uploads folder gets extracted,

@app.on_event("startup")
async def reindex_existing_files():
    existing = sorted(
        f for f in os.listdir(UPLOAD_FOLDER)
        if os.path.splitext(f)[1].lower() in SUPPORTED_EXTENSIONS
    )
    if not existing:
        return
    print(f"[startup] re-indexing {len(existing)} existing file(s)…")
    for fname in existing:
        if fname in file_names:
            continue
        path = os.path.join(UPLOAD_FOLDER, fname)
        try:
            text = extract_text(path)
            if not text.strip():
                continue
            file_names.append(fname)
            chunks = chunk_text(text, size=800, overlap=150)
            embs   = embed_model.encode(chunks)
            faiss_index.add(np.array(embs, dtype="float32"))
            documents.extend(chunks)
            for c in chunks:
                doc_meta.append({"text": c, "file_name": fname})
            print(f"[startup] indexed '{fname}': {len(chunks)} chunks")
        except Exception as e:
            print(f"[startup] skipping '{fname}': {e}")
    print(f"[startup] done — {len(documents)} total chunks ready")


# ── Pydantic models ───────────────────────────────────
class FilesBody(BaseModel):
    selected_files: Optional[List[str]] = []

class QuestionBody(BaseModel):
    question: str
    selected_files: Optional[List[str]] = []

class AnswerBody(BaseModel):
    session_id: str
    answer: str

class DeleteBody(BaseModel):
    filename: str


# ── Text extraction ───────────────────────────────────

def extract_text_pdf(path: str) -> str:
    doc = fitz.open(path)
    return "".join(page.get_text() for page in doc)

def _extract_table_text(table) -> list:
    parts = []
    for row in table.rows:
        row_texts = []
        for cell in row.cells:
            text = cell.text.strip()
            if text:
                row_texts.append(text)
            for nested_table in cell.tables:
                parts.extend(_extract_table_text(nested_table))
        if row_texts:
            parts.append(" | ".join(row_texts))
    return parts

def extract_text_docx(path: str) -> str:
    document = docx.Document(path)
    parts = []
    for para in document.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    for table in document.tables:
        parts.extend(_extract_table_text(table))
    return "\n".join(parts)

def extract_text_pptx(path: str) -> str:
    prs = Presentation(path)
    parts = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        parts.append(f"[Slide {slide_num}]")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    parts.append(text)
    return "\n".join(parts)

def extract_text_txt(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        with open(path, "r", encoding="latin-1") as f:
            return f.read()

def extract_text(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    extractors = {
        ".pdf":  extract_text_pdf,
        ".docx": extract_text_docx,
        ".pptx": extract_text_pptx,
        ".txt":  extract_text_txt,
    }
    if ext not in extractors:
        raise ValueError(f"Unsupported file type: {ext}")
    text = extractors[ext](path)
    print(f"[extract_text] {os.path.basename(path)} ({ext}): {len(text)} chars")
    return text


# ── Chunking ──────────────────────────────────────────

def chunk_text(text: str, size: int = 800, overlap: int = 150) -> list:
    """
    Sentence-aware chunking:
    Splits on sentence endings (.!?) or paragraph breaks so chunks
    are always complete thoughts, never cut mid-sentence.
    """
    units = re.split(r'(?<=[.!?])\s+|\n\n+', text)
    units = [u.strip() for u in units if u.strip()]

    chunks  = []
    current = ""

    for unit in units:
        if len(unit) > size:
            if current:
                chunks.append(current.strip())
                current = ""
            start = 0
            while start < len(unit):
                chunks.append(unit[start:start + size])
                start += size - overlap
            continue

        if len(current) + len(unit) + 1 <= size:
            current = (current + " " + unit).strip()
        else:
            if current:
                chunks.append(current.strip())
            overlap_text = current[-overlap:] if len(current) > overlap else current
            current = (overlap_text + " " + unit).strip()

    if current:
        chunks.append(current.strip())

    print(f"[chunk_text] produced {len(chunks)} sentence-aware chunks")
    return chunks


# ── Retrieval intent detection ────────────────────────

EXACT_KEYWORDS = [
    "show me", "show the", "display", "print", "output",
    "as it is", "as-is", "verbatim", "exact", "exactly",
    "word for word", "copy", "paste", "reproduce",
    "what does it say", "what does the document say",
    "read out", "read me", "give me the text",
    "extract", "pull out", "retrieve", "fetch",
    "what is written", "what is stated", "quote",
]

def detect_exact_request(question: str) -> bool:
    q = question.lower()
    return any(kw in q for kw in EXACT_KEYWORDS)

def detect_target_file(question: str, available_files: list):
    q_lower = question.lower()
    for fname in sorted(available_files, key=len, reverse=True):
        base = os.path.splitext(fname)[0].lower()
        if fname.lower() in q_lower or base in q_lower:
            print(f"[detect_target_file] matched '{fname}' in question")
            return fname
    return None

def detect_detail_request(question: str) -> bool:
    detail_keywords = [
        "explain", "elaborate", "detail", "in depth", "in-depth",
        "thoroughly", "comprehensive", "tell me more", "expand",
        "describe", "walk me through", "break down", "deep dive",
        "step by step", "step-by-step", "fully", "complete",
        "long", "everything about", "all about"
    ]
    return any(kw in question.lower() for kw in detail_keywords)


# ── Groq call ─────────────────────────────────────────

def call_groq(prompt: str, max_tokens: int = 4096) -> str:
    try:
        r = requests.post(
            GROQ_URL,
            headers={
                "Authorization": f"Bearer {GROQ_API_KEY}",
                "Content-Type":  "application/json"
            },
            json={
                "model":      GROQ_MODEL,
                "messages":   [{"role": "user", "content": prompt}],
                "max_tokens": max_tokens,
            },
            timeout=60,
        )
        resp = r.json()
        if "error" in resp:
            print(f"[call_groq] API error: {resp['error']}")
            return f"__GROQ_ERROR__: {resp['error'].get('message', str(resp['error']))}"
        if "choices" not in resp:
            print(f"[call_groq] Unexpected response: {resp}")
            return "__GROQ_ERROR__: No choices in response"
        return resp["choices"][0]["message"]["content"]
    except Exception as e:
        print(f"[call_groq] Exception: {e}")
        return f"__GROQ_ERROR__: {e}"


# ── Retrieval helpers ─────────────────────────────────

def chunks_for(selected_files: list):
    selected_files = selected_files or []
    if not selected_files:
        return documents[:], list(range(len(documents)))
    texts, idxs = [], []
    for i, m in enumerate(doc_meta):
        if m["file_name"] in selected_files:
            texts.append(m["text"]); idxs.append(i)
    return texts, idxs

def cosine_top_k(query_vec, texts, k: int = 5):
    if not texts:
        return []
    embs = embed_model.encode(texts)
    qv   = query_vec / (np.linalg.norm(query_vec) + 1e-9)
    sims = embs @ qv / (np.linalg.norm(embs, axis=1) + 1e-9)
    top  = np.argsort(sims)[::-1][:k]
    return [texts[i] for i in top]

def cosine_top_k_per_file(query_vec, selected_files: list, k_per_file: int = 3):
    selected_files = selected_files or []
    by_file: dict[str, list[str]] = {}
    for m in doc_meta:
        if m["file_name"] in selected_files:
            by_file.setdefault(m["file_name"], []).append(m["text"])
    all_retrieved = []
    for fname, chunks in by_file.items():
        if not chunks:
            continue
        top_chunks = cosine_top_k(query_vec, chunks, k=min(k_per_file, len(chunks)))
        for chunk in top_chunks:
            all_retrieved.append(f"[From: {fname}]\n{chunk}")
    return all_retrieved

def sample_chunks(selected_files: list, total: int = 12, labeled: bool = False) -> str:
    selected_files = selected_files or []
    files_to_use = (
        selected_files if selected_files
        else list(dict.fromkeys(m["file_name"] for m in doc_meta))
    )
    print(f"[sample_chunks] files_to_use={files_to_use}")
    by_file: dict[str, list[str]] = {}
    for m in doc_meta:
        if m["file_name"] in files_to_use:
            by_file.setdefault(m["file_name"], []).append(m["text"])
    print(f"[sample_chunks] chunk counts: { {k: len(v) for k, v in by_file.items()} }")
    if not by_file:
        return ""
    num_files    = len(by_file)
    per_file     = max(5, total // num_files)
    scaled_total = min(per_file * num_files, 60)
    sampled = []
    for fname, chunks in by_file.items():
        if len(chunks) <= per_file:
            selected = chunks[:]
        else:
            step     = len(chunks) / per_file
            indices  = [int(i * step) for i in range(per_file)]
            selected = [chunks[i] for i in indices]
        if labeled:
            block  = f"\n--- START OF DOCUMENT: {fname} ---\n"
            block += "\n".join(selected)
            block += f"\n--- END OF DOCUMENT: {fname} ---\n"
            sampled.append(block)
        else:
            sampled.extend(selected)
    if labeled:
        result = "\n\n".join(sampled)
    else:
        result = "\n\n".join(sampled[:scaled_total])
    print(f"[sample_chunks] final context: {len(result)} chars")
    return result

def trim_context_to_limit(context: str, char_limit: int = 7000) -> str:
    if len(context) <= char_limit:
        return context
    blocks = re.split(r"(?=--- START OF DOCUMENT:)", context)
    blocks = [b.strip() for b in blocks if b.strip()]
    if len(blocks) <= 1:
        trimmed = context[:char_limit]
        print(f"[trim_context] simple trim: {len(context)} → {len(trimmed)} chars")
        return trimmed
    per_block = char_limit // len(blocks)
    trimmed_blocks = []
    for block in blocks:
        if len(block) > per_block:
            lines        = block.split("\n")
            header       = lines[0] + "\n"
            body         = "\n".join(lines[1:])
            trimmed_body = body[:per_block - len(header)]
            trimmed_blocks.append(header + trimmed_body)
        else:
            trimmed_blocks.append(block)
    result = "\n\n".join(trimmed_blocks)
    print(f"[trim_context] proportional trim: {len(context)} → {len(result)} chars")
    return result

def extract_json_array(raw: str):
    cleaned = re.sub(r"```(?:json)?", "", raw).strip()
    start = cleaned.find("[")
    if start == -1:
        raise ValueError("No JSON array found")
    depth, i = 0, start
    while i < len(cleaned):
        if   cleaned[i] == "[": depth += 1
        elif cleaned[i] == "]": depth -= 1
        if depth == 0:
            return json.loads(cleaned[start:i+1])
        i += 1
    raise ValueError("Unterminated JSON array")


# ── Prompt builders ───────────────────────────────────

def build_answer_prompt(question: str, context: str, is_detail: bool) -> str:
    source_note = (
        "Each chunk of context is labeled with [From: filename] "
        "so you know which document it came from. "
        "If the answer is only in one document, mention which one.\n\n"
    )
    if is_detail:
        return f"""You are a knowledgeable study assistant. The user has explicitly asked for a DETAILED, IN-DEPTH explanation.

{source_note}Your answer MUST:
- Be thorough and comprehensive — aim for at least 300-500 words
- Cover all relevant aspects found in the context
- Use clear structure: headings, bullet points, or numbered steps where appropriate
- Explain the "why" and "how", not just the "what"
- Include examples or elaborations if the context supports it
- End with a brief summary of the key takeaway

Use ONLY the information from the context below. If the context does not have enough detail say so clearly, but still explain everything the context does contain.

Context:
{context}

Question: {question}

Detailed Answer:"""
    else:
        return f"""You are a helpful study assistant. Answer the question clearly and accurately based on the context below.

{source_note}- If the question is simple, a concise answer (2-4 sentences) is fine.
- If the question is naturally complex, give a fuller explanation.
- Do NOT pad or repeat yourself.
- Answer ONLY from the context provided.
- If the answer is found in a specific document, mention which one.

Context:
{context}

Question: {question}

Answer:"""

def build_exact_retrieval_prompt(question: str, context: str, source_file: str = None) -> str:
    source_hint = (
        f"The user is asking specifically about: {source_file}.\n"
        f"Focus ONLY on content from that document.\n\n"
        if source_file else ""
    )
    return f"""The user wants to see the EXACT text from the document.
Do NOT paraphrase, summarise, rephrase, or rewrite anything.

{source_hint}Instructions:
1. Find the passage(s) in the context below that best match the request.
2. Copy those passage(s) VERBATIM — word for word, character for character.
3. Present the quoted text in a clearly labelled block:

📄 Exact text from [{source_file or 'document'}]:
\"\"\"
[paste the verbatim passage here — do not change a single word]
\"\"\"

4. After the quote, add ONE short sentence stating where in the document this appears.
5. If the requested passage is not present in the context at all, say so clearly.

Context:
{context}

User request: {question}

Verbatim passage:"""


# ── Prune stale quiz sessions ─────────────────────────

def prune_stale_sessions():
    now   = time.time()
    stale = [k for k, v in quiz_sessions.items()
             if now - v.get("created_at", now) > QUIZ_SESSION_TTL]
    for k in stale:
        del quiz_sessions[k]
    if stale:
        print(f"[prune_sessions] Removed {len(stale)} stale session(s): {stale}")


# ══════════════════════════════════════════════════════
# UPLOAD — auto-replaces existing file instead of
# blocking re-uploads. Edit your Word doc locally,
# re-upload the same filename, and the new content
# is immediately available for Q&A / summary / quiz.
# ══════════════════════════════════════════════════════
@app.post("/upload")
async def upload(file: UploadFile = File(...)):
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return {"error": f"Unsupported file type '{ext}'. Supported: PDF, DOCX, PPTX, TXT"}

    path = os.path.join(UPLOAD_FOLDER, file.filename)

    # ── AUTO-REPLACE ──────────────────────────────────
    # If the same filename already exists, wipe its old
    # chunks from memory and rebuild FAISS before
    # indexing the new version. This means you can edit
    # a Word doc locally, re-upload it, and the updated
    # content is instantly reflected — no manual delete.
    if file.filename in file_names:
        print(f"[upload] '{file.filename}' exists — replacing with updated version...")

        global documents, faiss_index

        surviving_meta = [m for m in doc_meta if m["file_name"] != file.filename]
        surviving_docs = [m["text"] for m in surviving_meta]

        doc_meta.clear();  doc_meta.extend(surviving_meta)
        documents.clear(); documents.extend(surviving_docs)
        file_names.remove(file.filename)

        # Rebuild FAISS without the old file's vectors
        faiss_index = faiss.IndexFlatL2(384)
        if surviving_docs:
            embs = embed_model.encode(surviving_docs)
            faiss_index.add(np.array(embs, dtype="float32"))

        print(f"[upload] old chunks cleared. Remaining chunks: {len(documents)}")

    # ── Save new file to disk ─────────────────────────
    with open(path, "wb") as buf:
        shutil.copyfileobj(file.file, buf)

    # ── Extract text ──────────────────────────────────
    try:
        text = extract_text(path)
    except Exception as e:
        os.remove(path)
        return {"error": f"Could not extract text: {e}"}

    if not text.strip():
        os.remove(path)
        return {"error": "No readable text found. Make sure the document is not blank."}

    # ── Chunk → embed → index ─────────────────────────
    file_names.append(file.filename)
    chunks = chunk_text(text, size=800, overlap=150)
    embs   = embed_model.encode(chunks)

    faiss_index.add(np.array(embs, dtype="float32"))
    documents.extend(chunks)
    for c in chunks:
        doc_meta.append({"text": c, "file_name": file.filename})

    print(f"[upload] '{file.filename}' ({ext}): {len(chunks)} chunks. Total: {len(doc_meta)}")
    return {"message": "Uploaded", "chunks": len(chunks), "files": file_names}


# ── Debug ─────────────────────────────────────────────
@app.get("/debug")
async def debug():
    files_in_meta = {}
    for m in doc_meta:
        files_in_meta[m["file_name"]] = files_in_meta.get(m["file_name"], 0) + 1
    return {
        "file_names":      file_names,
        "doc_meta_counts": files_in_meta,
        "total_chunks":    len(documents),
        "total_meta":      len(doc_meta),
        "active_sessions": len(quiz_sessions),
    }


# ── Files ─────────────────────────────────────────────
@app.get("/files")
async def get_files():
    return {"files": file_names}


# ── Delete ────────────────────────────────────────────
@app.post("/delete")
async def delete_file(body: DeleteBody):
    fname = body.filename.strip()

    if fname not in file_names:
        return {"error": f"'{fname}' not found in uploaded files."}

    path = os.path.join(UPLOAD_FOLDER, fname)
    if os.path.exists(path):
        os.remove(path)
        print(f"[delete] removed from disk: {path}")

    file_names.remove(fname)

    global documents, faiss_index
    surviving_meta = [m for m in doc_meta if m["file_name"] != fname]
    surviving_docs = [m["text"] for m in surviving_meta]

    doc_meta.clear();  doc_meta.extend(surviving_meta)
    documents.clear(); documents.extend(surviving_docs)

    faiss_index = faiss.IndexFlatL2(384)
    if surviving_docs:
        embs = embed_model.encode(surviving_docs)
        faiss_index.add(np.array(embs, dtype="float32"))

    print(f"[delete] '{fname}' removed. Files left: {file_names}. Chunks: {len(documents)}")
    return {"status": "ok", "files": file_names}


# ── Query ─────────────────────────────────────────────
@app.post("/query")
async def query(body: QuestionBody):
    if not documents:
        return {"answer": "⚠️ Please upload a document first."}

    selected_files = body.selected_files or []

    target_file = detect_target_file(body.question, file_names)
    if target_file:
        if not selected_files or target_file in selected_files:
            print(f"[query] locking retrieval to: {target_file}")
            selected_files = [target_file]

    is_exact  = detect_exact_request(body.question)
    is_detail = detect_detail_request(body.question) and not is_exact
    print(f"[query] exact={is_exact} | detail={is_detail} | '{body.question}'")

    qv = embed_model.encode([body.question])[0]

    if selected_files and len(selected_files) > 1:
        k_per_file = 5 if is_exact else (4 if is_detail else 3)
        retrieved  = cosine_top_k_per_file(qv, selected_files, k_per_file=k_per_file)
        if not retrieved:
            return {"answer": "⚠️ No content found for the selected file(s)."}
        context = "\n\n".join(retrieved)
        print(f"[query] per-file retrieval: {len(retrieved)} chunks from {len(selected_files)} files")
    else:
        texts, _ = chunks_for(selected_files)
        if not texts:
            return {"answer": "⚠️ No content found for the selected file(s)."}
        k       = 8 if is_exact else (7 if is_detail else 5)
        context = "\n\n".join(cosine_top_k(qv, texts, k=k))
        print(f"[query] single-file retrieval: top-{k} chunks")

    if is_exact:
        single_source = target_file or (selected_files[0] if len(selected_files) == 1 else None)
        prompt = build_exact_retrieval_prompt(body.question, context, source_file=single_source)
    else:
        prompt = build_answer_prompt(body.question, context, is_detail)

    answer = call_groq(prompt)

    if answer.startswith("__GROQ_ERROR__"):
        return {"answer": "⚠️ AI service error. Please try again."}

    return {"answer": answer}


# Route alias — Flask frontend calls /ask; FastAPI exposes /query
@app.post("/ask")
async def ask(body: QuestionBody):
    return await query(body)


# ── Summary ───────────────────────────────────────────
@app.post("/summary")
async def summary(body: FilesBody):
    selected_files = body.selected_files or []
    print(f"[summary] selected_files={selected_files}")

    if not documents:
        return {"summary": "⚠️ Please upload a document first."}

    num_files = (
        len(selected_files) if selected_files
        else len(set(m["file_name"] for m in doc_meta))
    )

    context_raw = sample_chunks(selected_files, total=10 * num_files, labeled=True)
    if not context_raw:
        return {"summary": "⚠️ No content found for the selected file(s)."}

    context = trim_context_to_limit(context_raw, char_limit=8000)
    print(f"[summary] {len(context)} chars → Groq, {num_files} file(s)")

    if num_files > 1:
        prompt = f"""You are summarising {num_files} separate documents provided below.
Each document is clearly marked with --- START OF DOCUMENT: filename --- and --- END OF DOCUMENT: filename ---.

You MUST write a separate detailed summary section for EVERY document. Do not skip any.
Use the exact filename as the heading for each section.

Use this format:

📄 **Summary: [filename of document 1]**
[4-6 sentences summarising document 1 in detail]

📄 **Summary: [filename of document 2]**
[4-6 sentences summarising document 2 in detail]

(exactly {num_files} 📄 sections — one per document)

📌 **Combined Overview**
[3-4 sentences linking themes across all documents]

🔑 **Key Points Across All Documents**
- Point 1 — [1-2 sentence explanation]
- Point 2 — [1-2 sentence explanation]
- Point 3 — [1-2 sentence explanation]
- Point 4 — [1-2 sentence explanation]
- Point 5 — [1-2 sentence explanation]

📚 **Important Concepts**
[Explain 3-4 important concepts in detail, 2-3 sentences each]

Documents:
{context}"""
    else:
        prompt = f"""Write a thorough, well-structured summary. Be detailed — explain each point, do not just list it.

📌 **Overview**
[3-4 sentences about the main topic and purpose]

🔑 **Key Points**
- Point 1 — [1-2 sentence explanation]
- Point 2 — [1-2 sentence explanation]
- Point 3 — [1-2 sentence explanation]
- Point 4 — [1-2 sentence explanation]
- Point 5 — [1-2 sentence explanation]

📚 **Important Concepts**
[Explain 3-4 important concepts in detail, 2-3 sentences each]

Text:
{context}"""

    result = call_groq(prompt)

    if result.startswith("__GROQ_ERROR__"):
        return {"summary": "⚠️ AI service error generating summary. Please try again."}

    return {"summary": result}


# ── Quiz helpers ──────────────────────────────────────

def generate_questions(selected_files: list):
    selected_files = selected_files or []
    print(f"[generate_questions] selected_files={selected_files}")

    num_files       = len(selected_files) if selected_files else 1
    chunks_per_file = max(3, min(5, 30 // num_files))
    context_raw     = sample_chunks(
        selected_files,
        total=chunks_per_file * num_files,
        labeled=True
    )

    if not context_raw:
        return None

    context = trim_context_to_limit(context_raw, char_limit=6000)
    print(f"[generate_questions] context={len(context)} chars")

    questions_per_file = max(1, 10 // num_files)
    remainder          = 10 - (questions_per_file * num_files)

    def build_quiz_prompt(ctx: str) -> str:
        return f"""Generate exactly 10 multiple-choice questions from the {num_files} document(s) below.
Each document is marked with --- START OF DOCUMENT: filename --- and --- END OF DOCUMENT: filename ---.

IMPORTANT: Distribute questions evenly — {questions_per_file} per document{f', {remainder} extra from any document' if remainder else ''}.

Rules:
- Each question must have exactly 4 options.
- The "answer" field must EXACTLY match one of the strings in "options".
- Include a 1-2 sentence "explanation" for the correct answer.
- Return ONLY a valid JSON array — no markdown, no prose, no code fences,
  nothing before or after the array.

Format:
[
  {{
    "id": 1,
    "question": "...",
    "options": ["...", "...", "...", "..."],
    "answer": "...",
    "explanation": "..."
  }}
]

Documents:
{ctx}"""

    def try_generate(ctx: str):
        raw = call_groq(build_quiz_prompt(ctx), max_tokens=2048)
        if raw.startswith("__GROQ_ERROR__"):
            print(f"[generate_questions] Groq error: {raw}")
            return None
        try:
            questions = extract_json_array(raw)
            valid = []
            for q in questions:
                if all(k in q for k in ("question", "options", "answer", "explanation")):
                    if q["answer"] not in q["options"]:
                        for opt in q["options"]:
                            if opt.strip().lower() == q["answer"].strip().lower():
                                q["answer"] = opt
                                break
                    valid.append(q)
            print(f"[generate_questions] valid questions: {len(valid)}")
            return valid if valid else None
        except Exception as e:
            print(f"[generate_questions] parse failed: {e}\nRaw (first 400): {raw[:400]}")
            return None

    result = try_generate(context)
    if result:
        return result

    print("[generate_questions] Retrying with reduced context...")
    reduced = trim_context_to_limit(context_raw, char_limit=3000)
    return try_generate(reduced)


# ── Quiz: start ───────────────────────────────────────
@app.post("/quiz/start")
async def quiz_start(body: FilesBody):
    global quiz_counter
    prune_stale_sessions()

    if not documents:
        return {"error": "⚠️ Please upload a document first."}

    selected_files = body.selected_files or []
    if not selected_files:
        return {"error": "⚠️ Please select at least one document before starting a quiz."}

    questions = generate_questions(selected_files)
    if not questions:
        return {"error": "❌ Could not generate quiz questions. Please try again."}

    quiz_counter += 1
    sid = f"quiz_{quiz_counter}"

    quiz_sessions[sid] = {
        "questions":     questions,
        "current_index": 0,
        "score":         0,
        "created_at":    time.time(),
    }

    q = questions[0]
    return {
        "session_id":      sid,
        "question_number": 1,
        "total_questions": len(questions),
        "question":        q["question"],
        "options":         q["options"],
    }


# ── Quiz: answer ──────────────────────────────────────
@app.post("/quiz/answer")
async def quiz_answer(body: AnswerBody):
    sid = body.session_id
    if sid not in quiz_sessions:
        return {"error": "❌ Quiz session not found. Please start a new quiz."}

    session = quiz_sessions[sid]

    if time.time() - session.get("created_at", 0) > QUIZ_SESSION_TTL:
        del quiz_sessions[sid]
        return {"error": "❌ Quiz session expired. Please start a new quiz."}

    idx     = session["current_index"]
    total   = len(session["questions"])
    current = session["questions"][idx]

    is_correct = (body.answer.strip() == current["answer"].strip())
    if is_correct:
        session["score"] += 1

    feedback = {
        "is_correct":     is_correct,
        "correct_answer": current["answer"],
        "explanation":    current.get("explanation", ""),
        "your_answer":    body.answer,
    }

    if idx >= total - 1:
        feedback["quiz_finished"]   = True
        feedback["final_score"]     = session["score"]
        feedback["total_questions"] = total
        feedback["percentage"]      = round((session["score"] / total) * 100, 1)
        del quiz_sessions[sid]
        return feedback

    session["current_index"] += 1
    nq = session["questions"][session["current_index"]]
    feedback["next_question"] = {
        "question_number": session["current_index"] + 1,
        "total_questions": total,
        "current_score":   session["score"],
        "question":        nq["question"],
        "options":         nq["options"],
    }
    return feedback